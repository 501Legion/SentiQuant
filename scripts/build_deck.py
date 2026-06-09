"""community-opinion-agent 발표용 .pptx 생성기.
실행: venv/Scripts/python.exe scripts/build_deck.py
출력: docs/presentation/community-opinion-agent.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

KFONT = "맑은 고딕"
NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x2E, 0x5B, 0xFF)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
GREEN = RGBColor(0x1B, 0x8A, 0x5A)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_font(run, size, bold=False, color=GRAY, font=KFONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bar(slide, color=BLUE, h=0.18):
    s = slide.shapes.add_shape(1, 0, 0, SW, Inches(h))
    _fill(s, color)


def title_slide(title, subtitle, tag):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH); _fill(bg, NAVY)
    accent = s.shapes.add_shape(1, 0, Inches(5.0), SW, Inches(0.12)); _fill(accent, BLUE)
    tb = _box(s, 0.9, 2.2, 11.5, 2.0); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title; _set_font(r, 40, True, WHITE)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle; _set_font(r2, 20, False, RGBColor(0xC8,0xD2,0xE8))
    tg = _box(s, 0.9, 1.4, 11, 0.6); r = tg.text_frame.paragraphs[0].add_run()
    r.text = tag; _set_font(r, 14, True, BLUE)
    ft = _box(s, 0.9, 6.5, 11.5, 0.6); r = ft.text_frame.paragraphs[0].add_run()
    r.text = "auto_stock · branch: community-opinion-agent · 2026-05"; _set_font(r, 12, False, RGBColor(0x90,0x9A,0xB5))
    return s


def section_slide(num, title):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH); _fill(bg, NAVY)
    n = _box(s, 0.9, 2.6, 3, 1.5); r = n.text_frame.paragraphs[0].add_run()
    r.text = num; _set_font(r, 60, True, BLUE)
    tb = _box(s, 0.9, 3.7, 11.5, 1.5); r = tb.text_frame.paragraphs[0].add_run()
    r.text = title; _set_font(r, 32, True, WHITE)
    return s


def content_slide(title, bullets, foot=None):
    """bullets: list of (level, text, bold?) — level 0/1."""
    s = prs.slides.add_slide(BLANK)
    bar(s)
    tb = _box(s, 0.7, 0.45, 12, 1.0); r = tb.text_frame.paragraphs[0].add_run()
    r.text = title; _set_font(r, 28, True, NAVY)
    body = _box(s, 0.85, 1.7, 11.7, 5.2); tf = body.text_frame; tf.word_wrap = True
    first = True
    for item in bullets:
        lvl, text = item[0], item[1]
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else (NAVY if bold else GRAY)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        prefix = "" if lvl == 0 and bold else ("■  " if lvl == 0 else "–  ")
        r = p.add_run(); r.text = prefix + text
        _set_font(r, 18 if lvl == 0 else 15, bold, color)
        p.space_after = Pt(6)
    if foot:
        fb = _box(s, 0.85, 6.95, 11.7, 0.45); r = fb.text_frame.paragraphs[0].add_run()
        r.text = foot; _set_font(r, 12, False, GRAY)
    return s


def table_slide(title, headers, rows, foot=None, col_widths=None, highlight_col=None):
    s = prs.slides.add_slide(BLANK)
    bar(s)
    tb = _box(s, 0.7, 0.45, 12, 1.0); r = tb.text_frame.paragraphs[0].add_run()
    r.text = title; _set_font(r, 28, True, NAVY)
    nrows, ncols = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.7), Inches(1.7),
                            Inches(11.9), Inches(0.5 + 0.45 * len(rows)))
    table = gt.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = table.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rr = p.add_run(); rr.text = h; _set_font(rr, 13, True, WHITE)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            rr = p.add_run(); rr.text = str(val)
            col = NAVY if j == 0 else GRAY
            _set_font(rr, 12, j == 0, col)
    if foot:
        fb = _box(s, 0.7, 6.95, 12, 0.45); r = fb.text_frame.paragraphs[0].add_run()
        r.text = foot; _set_font(r, 12, False, GRAY)
    return s


# ============================================================ 슬라이드 구성
# 1. 타이틀
title_slide(
    "Community Opinion Agent",
    "커뮤니티 여론 트렌드 기반 트레이딩 의사결정 에이전트 (v0~v3)",
    "AI NATIVE QUANT · PDCA")

# 2. 한 장 요약
content_slide("한 장 요약 (TL;DR)", [
    (0, "무엇을: Reddit/WSB 커뮤니티 '여론'을 수집·해석해 매매를 판단하는 에이전트", True),
    (1, "급등 모멘텀 추격이 아니라 — 의견의 방향·지속성·합의도·노이즈·관심도 변화에 베팅"),
    (0, "어떻게: rule 신호(WSB V3)를 1차 후보로, 그 위에 필터·메모리·라우터를 얹음", True),
    (1, "Universe/Cost 게이팅 → 글품질/티커 필터 → 의견 스냅샷 → 메모리/리플렉션 → Decision Router"),
    (0, "결과: 신규 5모듈 + 테스트 86건 통과, Match Rate 98%, equal 회귀 0", True, GREEN),
    (1, "LLM은 자율 매매자가 아니라 '도구 결과를 해석·설명하는 보조 라우터' (기본 OFF)"),
], foot="branch: community-opinion-agent · 기존 trend-sizing 위 확장")

# 3. 문제 정의
content_slide("1. 문제 — 기존 여론 전략의 한계", [
    (0, "거래 대상 무차별", True), (1, "저유동·OTC·penny 종목까지 거래 → 왕복 수수료가 수익 잠식"),
    (0, "비용 미고려", True), (1, "gross(수수료 전)만 봄 → 실제 net 수익과 괴리"),
    (0, "글 품질·티커 오탐", True), (1, "Meme 글과 DD(심층분석) 글이 동일 취급, 'ALL/IT/NOW' 같은 일반어가 티커로 오인"),
    (0, "학습·근거 부재", True), (1, "과거 유사 사례를 기억 못 함, '왜 이 거래를 했는지' 구조화된 기록 없음"),
    (0, "NEW_SPIKE 과대평가 위험", True, RED), (1, "단발 언급 폭증을 강한 매수로 착각 → 급등추격화"),
])

# 4. 전략 정의
content_slide("2. 전략 정의 — 무엇에 베팅하는가", [
    (0, "이 전략은 '급등 모멘텀 추격'이 아니다", True, RED),
    (0, "커뮤니티 여론의 5가지 차원을 본다", True, NAVY),
    (1, "방향성 — 긍정/부정 (consensus_ratio)"),
    (1, "지속성 — 며칠 유지되는가 (persistence_days)"),
    (1, "합의도 — 한쪽으로 모이는가 (weighted bull/bear)"),
    (1, "노이즈 — 중립 비율이 낮은가 (neutral_ratio)"),
    (1, "관심도 변화 — 안정적 증가 vs 단발 폭증 (velocity_state)"),
    (0, "강한 매수 = 지속성 + 합의 + 낮은 노이즈 + 과거 성공이 동반될 때만", True, GREEN),
])

# 5. 아키텍처
content_slide("3. 시스템 아키텍처 (Option C — 독립 모듈 + 오케스트레이션)", [
    (0, "WSB V3 신호(run_pipeline) → Top-N 후보  [신호 엔진 불변]", True),
    (0, "매수 직전 에이전트 게이팅 파이프라인", True, NAVY),
    (1, "① UniverseFilter.decide()  → 거래 가능성·tier·size_multiplier"),
    (1, "② CostAwareTradeFilter.evaluate()  → 왕복비용 vs 기대 edge"),
    (1, "③ CommunityMemoryStore.retrieve()  → 유사 과거 사례"),
    (1, "④ DecisionRouter.decide()  → BUY/HOLD/SELL/REDUCE/SKIP/EXIT + 근거"),
    (0, "action==BUY만 매수 · 9-factor OpinionTrendSizer로 사이징", True),
    (0, "청산: 기존 5단계 유지 + opinion_reversal · 종료 후 snapshot/reflection 저장", True),
], foot="reddit_backtester가 오케스트레이션 · signals.py/backtester.py/뉴스 경로 불가침")

# 6. v0 Universe
table_slide("4. v0 — Universe Filter (거래 대상 선별)",
    ["tier", "정의", "size 배수"],
    [["CORE", "S&P500 ∪ Nasdaq100 인덱스 대형주", "1.0"],
     ["EXPANDED", "인덱스 외 대형/중형 + 유동성·시총 통과", "0.5"],
     ["COMMUNITY_LIQUID", "인덱스 외지만 커뮤니티 관심 + 유동성 통과", "0.5"],
     ["BLOCKED", "저유동·OTC·penny·티커오탐·시총 미달", "—"]],
    foot="6 모드: sp500_only ⊂ liquid_us ⊂ community_liquid(기본) · custom_watchlist 등 · 정적 JSON + OHLCV 유동성",
    col_widths=[2.6, 7.3, 2.0])

# 7. v0 Cost
content_slide("4. v0 — Cost-aware Trade Filter (비용 대비 edge)", [
    (0, "왕복비용 = 수수료×2 + 슬리피지 + 스프레드 = 0.7%", True),
    (0, "기대 움직임(edge) proxy 우선순위", True, NAVY),
    (1, "① ATR% → ② 최근 변동폭 → ③ 의견 확신도(conviction)"),
    (0, "edge < 왕복비용 × 2.0 → SKIP (비용이 기대수익 잡아먹음)", True, RED),
    (1, "경계 구간 → DOWNSIZE (size factor 0.7)"),
    (0, "→ gross/net 동시 출력, 수수료·슬리피지·turnover 추적", True, GREEN),
], foot="replay는 ATR 부재 → OHLCV 변동성 proxy 사용 (라이브는 ATR)")

# 8. v1 품질/티커
content_slide("5. v1 — 글 품질 & 티커 오탐 필터", [
    (0, "Source Quality (flair 가중)", True, NAVY),
    (1, "DD 1.5 · News/Fundamentals 1.2 · Discussion 1.0 · Daily Thread 0.5 · Meme/Gain/Loss 0.0"),
    (1, "title 멘션(2.0) > body(1.0) > comment(0.5) 가중"),
    (0, "Ticker Ambiguity 필터", True, NAVY),
    (1, "'ALL/IT/NOW/COST...' 같은 일반어는 $ 접두사($ALL) 있을 때만 인정"),
    (1, "단일문자 티커는 $F 처럼 $ 있어야 인정 · 제외 건 통계 집계"),
    (0, "→ 가중 카운트로 weighted_bullish/bearish 계산 (DailyOpinionSnapshot)", True, GREEN),
])

# 9. v1 Sizer
content_slide("5. v1 — OpinionTrendSizer (9-factor 사이징)", [
    (0, "final_size_factor = clamp( 곱(9 factor), 0.0 ~ 1.3 )", True, NAVY),
    (1, "opinion_score × trend × persistence × consensus × neutral × attention"),
    (1, "× source_quality × universe_multiplier × cost_risk_factor  (신규 3개)"),
    (0, "진입 게이팅 (0주 = 진입 제외)", True),
    (1, "score<60  또는  neutral>0.70  또는  consensus<1.5  또는  cost SKIP"),
    (0, "NEW_SPIKE 단독 → 0.5배 축소 (급등추격 방지) · 최대 1.3배 제한", True, RED),
    (0, "신규 factor 데이터 없으면 1.0 → 기존 동작 회귀 0", True, GREEN),
])

# 10. v2 memory/reflection
content_slide("6. v2 — Community Memory & Reflection (학습)", [
    (0, "CommunityMemoryStore (MemoryBackend 추상 → Jsonl / InMemory)", True, NAVY),
    (1, "과거 의견 스냅샷·리플렉션 저장 + 유사 사례 검색(휴리스틱: symbol·tier·score·키워드)"),
    (1, "향후 Chroma/Faiss(vector DB)로 교체 가능하도록 인터페이스 분리"),
    (0, "LowLevelReflection — 의견 신호 → 이후 가격(1/3/7/14일 수익률)", True),
    (1, "result_label: success_1d/3d/7d · delayed · noisy · failed"),
    (0, "HighLevelReflection — 실제 매매 entry/exit 분석", True),
    (1, "net_pnl·cost_drag·decision_quality(good/bad_entry·risk_mgmt 등)·lesson"),
], foot="백테스트는 run-local 메모리로 결정성 유지 (전역 파일 미조회)")

# 11. v3 router
content_slide("7. v3 — Decision Router (rule 기본 + LLM 보조)", [
    (0, "LLM은 자율 매매자가 아니다 — 도구 결과 해석·승인/축소/보류 라우터", True, RED),
    (0, "rule-based (기본)", True, NAVY),
    (1, "BUY 승인: 신호 BUY/STRONG_BUY + 합의·지속성·낮은 노이즈 + universe/cost 통과"),
    (1, "SKIP/REDUCE/SELL/EXIT: 노이즈 급증·합의 붕괴·과거 실패 패턴 등"),
    (0, "LLM router (선택, 기본 OFF)", True, NAVY),
    (1, "strict JSON 출력 · 실호출은 gpt-5.4-mini · confidence 낮으면 rule 우선"),
    (0, "8개 하드 안전장치", True, GREEN),
    (1, "rule SKIP을 LLM이 BUY로 못 뒤집음 · neutral/consensus/ambiguity/universe/cost/cash/no-position 차단"),
])

# 12. 청산 + 회귀
content_slide("8. 리스크 관리 & 회귀 보호", [
    (0, "청산 5단계 (순서 유지)", True, NAVY),
    (1, "① opinion_reversal(의견역전) ② RSI 과매수 ③ Gap Down -5% ④ Stop-Loss -7% ⑤ Trailing Stop -5%"),
    (0, "opinion_reversal: neutral 급증 · 합의 붕괴 · score 역전 · 추세 하락 · bearish 급증", True),
    (0, "회귀 보호 — 최상위 제약", True, GREEN),
    (1, "에이전트 게이팅은 opinion_trend에서만 · equal은 절대 미게이팅"),
    (1, "모든 신규 필터 config flag → OFF 시 기존 동작 byte 동일"),
    (1, "scripts/regression_check_reddit.py: equal 결과 diff 시 exit 1"),
])

# 13. 백테스트 결과 (4버전)
table_slide("9. 백테스트 — 4버전 비교 (2026-05-17~25, finbert-wsb)",
    ["설정 (sizing/universe)", "청산거래", "net%", "gross%", "수수료$", "universe skip", "router"],
    [["equal / community_liquid", "0", "-0.09", "+0.00", "86", "0", "(미게이팅)"],
     ["opinion_trend / community_liquid", "0", "-0.03", "+0.00", "29", "2", "BUY 5 / SKIP 2"],
     ["opinion_trend / sp500_only", "0", "+0.00", "+0.00", "0", "7", "BUY 1 / SKIP 7"],
     ["opinion_trend / liquid_us", "0", "+0.00", "+0.00", "0", "7", "BUY 1 / SKIP 7"]],
    foot="핵심: universe 모드가 실제로 차별화(skip 2→7, BUY 5→1) · opinion_trend가 equal보다 선별적(수수료 86→29)",
    col_widths=[3.6, 1.4, 1.2, 1.2, 1.4, 1.7, 1.4])

# 14. 결과 해석
content_slide("9. 백테스트 — 무엇을 확인했나", [
    (0, "✅ 기능 검증 성공 (실데이터 FinBERT)", True, GREEN),
    (1, "universe 모드 차별화: community_liquid는 후보 통과 多, sp500_only는 비인덱스 7건 차단"),
    (1, "opinion_trend가 equal보다 선별적 매수 → 수수료 86$ → 29$"),
    (1, "라우터·skip·비용 metric 전부 실데이터에서 정상 집계"),
    (0, "⚠️ 성과(수익률)는 판단 보류 — 데이터 한계", True, RED),
    (1, "윈도우 8거래일뿐 → 진입 포지션이 청산 전 기간 종료 → 청산거래 0건"),
    (1, "net%는 미청산 평가손 + 수수료 드래그만 반영 → 수익성 결론 불가"),
    (0, "결론: 구조·기능은 검증 완료, 성과 검증은 30일+ 데이터 축적 후", True, NAVY),
])

# 15. 품질
table_slide("10. 품질 & 검증",
    ["항목", "결과"],
    [["단위 테스트", "신규 73 + 기존 13 = 86건 전부 통과"],
     ["Match Rate (설계 대비 구현)", "98% (Critical 0, Important 0)"],
     ["Success Criteria", "11 / 11 충족"],
     ["equal 회귀", "0 (regression_check exit 0, 결정성 검증)"],
     ["신규 코드", "5 모듈 + 7 테스트 + 시드데이터 + 스크립트"],
     ["불가침 준수", "signals.py · backtester.py · 뉴스 경로 무수정"]],
    foot="PDCA: PM─ → Plan → plan-plus → Design(Option C) → Do(10 모듈) → Check 94% → Act → Report 98%",
    col_widths=[5.0, 6.9])

# 16. 한계 & 향후
content_slide("11. 한계 & 향후 로드맵", [
    (0, "현재 한계", True, RED),
    (1, "데이터 8~17일로 부족 → 성과 통계 검정력 낮음 (1순위: 일일 수집)"),
    (1, "factor 수치는 임의 초기값 · universe 리스트는 시드(105/71)"),
    (0, "향후 (별도 PDCA 사이클)", True, NAVY),
    (1, "① Grid Search — 파라미터 자동 튜닝 (walk-forward 과최적화 방지)"),
    (1, "② 네이버 종토방 수집기 — 한국 주식 + 한국어 감성모델 (KIS 실거래는 이미 연동)"),
    (1, "③ Vector DB(Chroma/Faiss) 실연동 — 의미 기반 유사 사례 검색 (인터페이스 이미 분리)"),
    (1, "④ 3-브랜치 수익률 비교 + Streamlit 대시보드 (읽기전용 우선 권장)"),
])

# 17. 결론
content_slide("12. 결론", [
    (0, "커뮤니티 여론을 '비용·품질·지속성·학습' 관점에서 해석하는 의사결정 인프라 완성", True, NAVY),
    (1, "급등추격이 아니라 의견 지속성·합의도·비용 효율에 베팅"),
    (1, "LLM은 보조 라우터(설명자)일 뿐 — 안전장치로 통제, 기본 OFF"),
    (0, "구조·기능 검증 완료(98%, 회귀 0) · 성과 검증은 데이터 축적이 관건", True, GREEN),
    (0, "다음 한 걸음: 데이터 수집 자동화 → universe 모드별 net 수익률 실측", True, BLUE),
])

os.makedirs("docs/presentation", exist_ok=True)
out = "docs/presentation/community-opinion-agent.pptx"
prs.save(out)
print(f"saved: {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
